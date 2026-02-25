#!/usr/bin/env python3
#
# md-pptx-injector.py
#

from __future__ import annotations

import argparse
import logging
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

__version__ = "0.9.4.0"

# -------------------------
# Constants
# -------------------------
DEFAULT_INDENT_SPACES = 2
MAX_BULLET_LEVEL = 4
HEADING_LEVEL_PAGE_BREAK_THRESHOLD = 3
HEADING_LEVEL_BASE_OFFSET = 4
CODE_BLOCK_LEFT = Inches(0.5)
CODE_BLOCK_TOP = Inches(0.5)
CODE_BLOCK_WIDTH = Inches(9.0)
CODE_BLOCK_HEIGHT = Inches(4.0)

logger = logging.getLogger(__name__)

def setup_logging(verbose: bool) -> None:
    """Configure log level. -v enables DEBUG so all diagnostic output is visible.
    force=True ensures any existing handlers (e.g. added by python-pptx) are replaced.
    """
    level = logging.DEBUG if verbose else logging.WARNING
    logging.basicConfig(level=level, format='%(message)s', stream=sys.stdout, force=True)

def get_app_dir() -> Path:
    """Get application directory (supports frozen exe)."""
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent
    return Path(__file__).resolve().parent

def resolve_path(base_dir: Path, p: str, prefer_cwd: bool = True) -> Path:
    """Resolve path: prefer CWD for CLI usability, fall back to app dir."""
    pp = Path(p)
    if pp.is_absolute():
        return pp
    cwd_path = (Path.cwd() / pp).resolve()
    if prefer_cwd and cwd_path.exists():
        return cwd_path
    return (base_dir / pp).resolve()

# -------------------------
# Regex patterns
# FIX #2: RE_LAYOUT and RE_PLACEHOLDER both use --\s*> to tolerate optional space before >
# -------------------------
RE_LAYOUT = re.compile(r'<!--\s*layout\s*=\s*"([^"]+)"\s*--\s*>')
RE_NEW_PAGE = re.compile(r'<!--\s*new_page(?:\s*=\s*"([^"]+)")?\s*--\s*>')
RE_PLACEHOLDER = re.compile(r'<!--\s*placeholder\s*=\s*"([^"]+)"\s*--\s*>')
RE_HEADING = re.compile(r"^(#{1,6})\s+(.*)$")
RE_KV = re.compile(r"^\s*(title|subtitle|author|toc|toc_title|indent|font_size_l[0-4])\s*:\s*(.*)\s*$", re.IGNORECASE)
RE_BULLET = re.compile(r"^(\s*)([-*+])\s+(.*)$")
RE_NUMBERED = re.compile(r"^(\s*)(\d+)\.\s+(.*)$")
RE_IMAGE = re.compile(r"!\[(.*?)\]\((.*?)\)")
RE_HTML_TAGS = re.compile(r'<(/?)([biu])>', re.IGNORECASE)
RE_CODE_FENCE = re.compile(r"^```")
RE_FONT_SIZE_ATTRS = re.compile(r'font_size_l([0-4])\s*=\s*(\d+(?:\.\d+)?)')

# -------------------------
# Data classes
# -------------------------
@dataclass
class TextRun:
    """A fragment of text with inline formatting."""
    text: str
    bold: bool = False
    italic: bool = False
    underline: bool = False

@dataclass
class ParaSpec:
    """Specification of a paragraph."""
    text: str
    bullet: bool = False
    level: int = 0
    empty: bool = False
    runs: list[TextRun] = field(default_factory=list)

@dataclass
class TOCEntry:
    """A table-of-contents entry."""
    text: str
    level: int
    target_slide: Any

# -------------------------
# PPTX XML helpers
# -------------------------
def _get_or_add_pPr(paragraph):
    return paragraph._p.get_or_add_pPr()

def set_bullet_none(paragraph) -> None:
    try:
        p_pr = _get_or_add_pPr(paragraph)
        tags = ["a:buNone", "a:buChar", "a:buAutoNum", "a:buBlip",
                "a:buClr", "a:buSzPct", "a:buSzPts", "a:buFont"]
        for tag in tags:
            for el in p_pr.findall(qn(tag)):
                p_pr.remove(el)
        p_pr.insert(0, OxmlElement("a:buNone"))
    except Exception as e:
        logger.debug(f"set_bullet_none failed: {e}")

def clear_bullet_override(paragraph) -> None:
    try:
        p_pr = _get_or_add_pPr(paragraph)
        for el in p_pr.findall(qn("a:buNone")):
            p_pr.remove(el)
    except Exception as e:
        logger.debug(f"clear_bullet_override failed: {e}")

def send_to_back(shape) -> None:
    """Send shape to the back of the z-order.
    spTree structure requires <p:nvGrpSpPr> and <p:grpSpPr> as the first two
    children — inserting at index 0 would corrupt the XML.
    Instead, find the first actual shape child and insert before it.
    """
    try:
        sp = shape._element
        parent = sp.getparent()
        parent.remove(sp)
        # Find insertion point: after required metadata elements
        REQUIRED_TAGS = {
            "nvGrpSpPr", "grpSpPr",  # spTree required children
        }
        insert_idx = 0
        for i, child in enumerate(parent):
            local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if local in REQUIRED_TAGS:
                insert_idx = i + 1
        parent.insert(insert_idx, sp)
    except Exception as e:
        logger.debug(f"send_to_back failed: {e}")

def bring_to_front(shape) -> None:
    try:
        sp = shape._element
        parent = sp.getparent()
        parent.remove(sp)
        parent.append(sp)
    except Exception as e:
        logger.debug(f"bring_to_front failed: {e}")

def set_ea_font(run, typeface: str = "MS Gothic") -> None:
    """Set East Asian font using DrawingML a:ea element."""
    try:
        rPr = run._r.get_or_add_rPr()
        for el in rPr.findall(qn("a:ea")):
            rPr.remove(el)
        ea = OxmlElement("a:ea")
        ea.set("typeface", typeface)
        rPr.append(ea)
    except Exception as e:
        logger.debug(f"set_ea_font failed: {e}")

def add_slide_hyperlink(run, source_slide, target_slide) -> None:
    """Add a slide-to-slide hyperlink via hlinkClick XML element."""
    try:
        rel_type = (
            "http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships/slide"
        )
        rId = source_slide.part.relate_to(target_slide.part, rel_type)
        rPr = run._r.get_or_add_rPr()
        for el in rPr.findall(qn("a:hlinkClick")):
            rPr.remove(el)
        hlink = OxmlElement("a:hlinkClick")
        hlink.set(qn("r:id"), rId)
        hlink.set("action", "ppaction://hlinksldjump")
        rPr.append(hlink)
    except Exception as e:
        logger.debug(f"add_slide_hyperlink failed: {e}")

# -------------------------
# PPTX shape finders
# -------------------------
def find_layout_by_name(prs: Presentation, layout_name: str):
    """Find slide layout by name.
    FIX #1: returns None instead of raising ValueError when layout is not found.
    Callers are responsible for fallback handling.
    """
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                return layout
    return None

def find_title_shape(slide):
    try:
        if slide.shapes.title:
            return slide.shapes.title
    except Exception as e:
        logger.debug(f"shapes.title access failed: {e}")
    for shp in slide.shapes:
        if "title" in getattr(shp, "name", "").lower():
            return shp
        try:
            if shp.placeholder_format and shp.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE
            ):
                return shp
        except (AttributeError, KeyError) as e:
            logger.debug(f"Title placeholder check failed: {e}")
    return None

def find_body_text_shape_excluding(slide, exclude_ids: set[int]):
    skip = {
        PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
        PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.DATE,
        PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.HEADER,
        PP_PLACEHOLDER.SLIDE_NUMBER,
    }
    for shp in getattr(slide, "placeholders", []):
        if id(shp) in exclude_ids or not getattr(shp, "has_text_frame", False):
            continue
        try:
            p_format = getattr(shp, "placeholder_format", None)
            if p_format and p_format.type in skip:
                continue
        except (AttributeError, KeyError) as e:
            logger.debug(f"Placeholder skip check failed: {e}")
            continue
        return shp
    return None

def find_shape_by_name(slide, shape_name: str):
    for shp in slide.shapes:
        if getattr(shp, "name", None) == shape_name:
            return shp
    try:
        target_idx = None
        for lph in slide.slide_layout.placeholders:
            if getattr(lph, "name", None) == shape_name:
                target_idx = lph.placeholder_format.idx
                break
        if target_idx is not None:
            for sph in slide.placeholders:
                if sph.placeholder_format.idx == target_idx:
                    return sph
    except Exception as e:
        logger.debug(f"find_shape_by_name layout lookup failed: {e}")
    return None

def dump_slide_shapes(slide, page_no: int) -> None:
    """Debug: dump all shapes on a slide (visible with -v)."""
    logger.debug(f"  [shapes on slide]")
    for i, shp in enumerate(slide.shapes):
        name = getattr(shp, "name", "")
        has_tf = getattr(shp, "has_text_frame", False)
        is_ph = getattr(shp, "is_placeholder", False)
        ph_idx = ph_type = None
        if is_ph:
            try:
                ph_idx = shp.placeholder_format.idx
                ph_type = shp.placeholder_format.type
            except (AttributeError, KeyError):
                pass
        st = getattr(shp, "shape_type", None)
        logger.debug(
            f"    - #{i}: name={name!r}, shape_type={st}, "
            f"is_placeholder={is_ph}, ph_idx={ph_idx}, ph_type={ph_type}, "
            f"has_text_frame={has_tf}"
        )

# -------------------------
# Markdown parsing
# -------------------------
def parse_font_sizes_from_comment(text: str) -> dict[int, float]:
    """Extract font_size_lN=X attributes from an HTML comment string.
    Returns {level: pt} for each valid font_size_lN found.
    Example: '<!-- layout="hoge" font_size_l1=16 font_size_l2=14 -->'
             -> {1: 16.0, 2: 14.0}
    """
    result: dict[int, float] = {}
    for m in RE_FONT_SIZE_ATTRS.finditer(text):
        level = int(m.group(1))
        try:
            pt = float(m.group(2))
            if pt > 0:
                result[level] = pt
        except ValueError:
            pass
    return result
def parse_with_stack(text: str, verbose: bool, page_no: int) -> list[TextRun]:
    """Stack-based inline HTML tag parser (<b>, <i>, <u>).
    Falls back to a single plain TextRun on mismatched or unclosed tags.
    """
    runs, pos, stack = [], 0, []
    for m in RE_HTML_TAGS.finditer(text):
        content = text[pos:m.start()]
        if content:
            runs.append(TextRun(
                text=content,
                bold='b' in stack,
                italic='i' in stack,
                underline='u' in stack,
            ))
        is_closing = m.group(1) == '/'
        tag = m.group(2).lower()
        if is_closing:
            if not stack or stack[-1] != tag:
                if verbose:
                    logger.info(f"Slide {page_no}: Tag mismatch in '{text}'. Skipping formatting.")
                return [TextRun(text=text)]
            stack.pop()
        else:
            stack.append(tag)
        pos = m.end()
    remainder = text[pos:]
    if remainder:
        if stack:
            if verbose:
                logger.info(f"Slide {page_no}: Unclosed tag in '{text}'. Skipping formatting.")
            return [TextRun(text=text)]
        runs.append(TextRun(text=remainder))
    return runs if runs else [TextRun(text=text)]

def skip_tables_and_images(lines: list[str]) -> list[str]:
    out, i = [], 0
    while i < len(lines):
        ln = lines[i]
        if RE_IMAGE.search(ln.strip()):
            i += 1
            continue
        if ln.strip().startswith("|") and i + 1 < len(lines):
            sep = lines[i + 1].strip()
            if sep.startswith("|") and "-" in sep:
                i += 2
                while i < len(lines) and lines[i].strip():
                    i += 1
                continue
        out.append(ln)
        i += 1
    return out

def build_paragraphs_from_lines(
    lines: list[str], indent: int, verbose: bool, page_no: int
) -> tuple[list[ParaSpec], list[str]]:
    paras, code_blocks = [], []
    in_code, current_code = False, []

    content_lines = skip_tables_and_images(list(lines))
    # Trim leading/trailing blank lines to avoid spurious empty paragraphs
    while content_lines and not content_lines[0].strip():
        content_lines.pop(0)
    while content_lines and not content_lines[-1].strip():
        content_lines.pop()

    safe_indent = indent if indent > 0 else DEFAULT_INDENT_SPACES
    base_level = 0

    for raw in content_lines:
        line = raw.rstrip("\n")
        if RE_CODE_FENCE.match(line.strip()):
            if in_code:
                code_blocks.append("\n".join(current_code))
                current_code, in_code = [], False
            else:
                in_code = True
            continue
        if in_code:
            current_code.append(line)
            continue
        if not line.strip():
            paras.append(ParaSpec(text="", empty=True))
            continue

        h_m = RE_HEADING.match(line.strip())
        if h_m and len(h_m.group(1)) >= HEADING_LEVEL_BASE_OFFSET:
            base_level = len(h_m.group(1)) - HEADING_LEVEL_BASE_OFFSET
            txt = h_m.group(2).strip()
            paras.append(ParaSpec(
                text=txt, level=base_level,
                runs=parse_with_stack(txt, verbose, page_no),
            ))
            continue

        num_m = RE_NUMBERED.match(line)
        bul_m = RE_BULLET.match(line)
        if num_m or bul_m:
            match = num_m if num_m else bul_m
            idt = len(match.group(1).replace("\t", " " * safe_indent))
            lv = min(base_level + (idt // safe_indent), MAX_BULLET_LEVEL)
            txt = (
                f"{num_m.group(2)}. {num_m.group(3).strip()}"
                if num_m else bul_m.group(3).strip()
            )
            paras.append(ParaSpec(
                text=txt, bullet=True, level=lv,
                runs=parse_with_stack(txt, verbose, page_no),
            ))
            continue

        paras.append(ParaSpec(
            text=line.strip(), level=base_level,
            runs=parse_with_stack(line.strip(), verbose, page_no),
        ))

    if in_code and current_code:
        code_blocks.append("\n".join(current_code))

    return paras, code_blocks

def split_pages(md_text: str) -> list[list[str]]:
    """
    Split markdown into pages:
    - YAML front matter at document start: kept as part of first page.
      A level-1 heading (#) immediately following YAML is also kept on the
      same page (title slide pattern: YAML + # Title).
    - <!-- new_page --> / <!-- new_page="layout" -->: unconditional page break
    - <!-- layout="name" --> immediately followed by #/##/###: triggers page break
    - ##/### heading lines: trigger page break
    - # heading: triggers page break ONLY when YAML front matter is not present
      on the current page
    """
    lines = md_text.splitlines()
    pages: list[list[str]] = []
    cur: list[str] = []
    in_yaml = False
    yaml_in_cur = False  # True while current page contains YAML front matter
    i = 0

    def flush(new_cur: list[str] | None = None) -> None:
        nonlocal cur, yaml_in_cur
        if any(l.strip() for l in cur):
            pages.append(cur)
        cur = new_cur if new_cur is not None else []
        yaml_in_cur = False

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not any(l.strip() for l in cur) and stripped == "---" and not in_yaml:
            in_yaml = True
            yaml_in_cur = True
            cur.append(line)
            i += 1
            continue
        if in_yaml:
            cur.append(line)
            if stripped == "---":
                in_yaml = False
            i += 1
            continue

        if stripped == "---":
            i += 1
            continue

        if RE_NEW_PAGE.match(stripped):
            flush([line])
            i += 1
            continue

        m_layout = RE_LAYOUT.search(stripped)
        if m_layout:
            next_stripped = lines[i + 1].strip() if i + 1 < len(lines) else ""
            h_m = RE_HEADING.match(next_stripped)
            if h_m and len(h_m.group(1)) <= HEADING_LEVEL_PAGE_BREAK_THRESHOLD:
                flush([line, lines[i + 1]])
                i += 2
            else:
                logger.warning(
                    f"layout tag ignored (next line is not a heading): {stripped!r}"
                )
                i += 1
            continue

        h_m = RE_HEADING.match(stripped)
        if h_m:
            level = len(h_m.group(1))
            if level <= HEADING_LEVEL_PAGE_BREAK_THRESHOLD:
                # Level-1 heading (#): keep on same page when YAML is present
                if level == 1 and yaml_in_cur:
                    cur.append(line)
                    i += 1
                    continue
                flush([line])
                i += 1
                continue

        cur.append(line)
        i += 1

    if cur:
        pages.append(cur)
    return pages

def extract_front_matter(page_lines: list[str]) -> tuple[dict[str, Any], list[str]]:
    i = 0
    while i < len(page_lines) and not page_lines[i].strip():
        i += 1
    if i >= len(page_lines) or page_lines[i].strip() != "---":
        return {}, page_lines
    j = i + 1
    while j < len(page_lines):
        if page_lines[j].strip() == "---":
            return parse_simple_yaml(page_lines[i + 1:j]), page_lines[j + 1:]
        j += 1
    return {}, page_lines

def parse_simple_yaml(yaml_lines: list[str]) -> dict[str, Any]:
    data = {}
    for line in yaml_lines:
        m = RE_KV.match(line)
        if m:
            key, val = m.group(1).lower(), m.group(2).strip()
            if key == "toc":
                data[key] = val.lower() == "true"
            elif key == "indent":
                try:
                    data[key] = int(val)
                except ValueError:
                    logger.warning(f"Invalid indent value '{val}', using default.")
            elif key.startswith("font_size_l"):
                # font_size_l0 .. font_size_l4 -> stored as float pt values
                level = int(key[-1])  # last char is 0-4
                try:
                    pt = float(val)
                    if pt <= 0:
                        raise ValueError("must be positive")
                    data[key] = pt
                except ValueError:
                    logger.warning(f"Invalid {key} value '{val}', ignoring.")
            else:
                data[key] = val
    return data

# -------------------------
# Slide content helpers
# -------------------------
def set_text_lines(shp, lines: list[str]) -> None:
    if not shp or not getattr(shp, "has_text_frame", False):
        return
    tf = shp.text_frame
    tf.clear()
    if not lines:
        return
    for i, s in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = s
        set_bullet_none(p)

def parse_title_page_info(front: dict, lines: list[str]) -> dict:
    info = dict(front)
    for ln in lines:
        m = RE_HEADING.match(ln.strip())
        if m and len(m.group(1)) == 1:
            info["title"] = m.group(2).strip()
        kv_m = RE_KV.match(ln)
        if kv_m:
            info[kv_m.group(1).lower()] = kv_m.group(2).strip()
    return info

def parse_markdown_table(lines: list[str]):
    """Parse a Markdown table.
    Returns (headers, rows, col_widths, col_aligns) or None.
    col_widths: list of dash counts from the separator row (proportional widths)
    col_aligns: list of 'left'/'center'/'right' per column
    """
    tbl = [ln.strip() for ln in lines if ln.strip()]
    if len(tbl) < 2 or not (tbl[0].startswith("|") and tbl[1].startswith("|")):
        return None
    h = [c.strip() for c in tbl[0].strip("|").split("|")]
    sep_cells = [c.strip() for c in tbl[1].strip("|").split("|")]
    col_widths = []
    col_aligns = []
    for cell in sep_cells:
        dashes = cell.replace(":", "")
        col_widths.append(max(len(dashes), 1))
        if cell.startswith(":") and cell.endswith(":"):
            col_aligns.append("center")
        elif cell.endswith(":"):
            col_aligns.append("right")
        else:
            col_aligns.append("left")
    r = [
        [c.strip() for c in ln.strip("|").split("|")]
        for ln in tbl[2:] if ln.strip().startswith("|")
    ]
    return h, r, col_widths, col_aligns

def _write_run(paragraph, run_spec: TextRun) -> None:
    """FIX #3: write a single TextRun into a paragraph at the run level.
    Using run.font rather than paragraph.font avoids overriding template
    paragraph-level styles (size, color, etc.) while still applying bold/italic/underline.
    """
    r = paragraph.add_run()
    r.text = run_spec.text
    try:
        if run_spec.bold:
            r.font.bold = True
        if run_spec.italic:
            r.font.italic = True
        if run_spec.underline:
            r.font.underline = True
    except Exception as e:
        logger.debug(f"Font style apply failed: {e}")

def write_paragraphs_to_shape(
    shape, paras: list[ParaSpec], append: bool, blank_before_append: bool,
    font_sizes: dict[int, float] | None = None,
) -> None:
    """Write paragraphs to a text shape.

    font_sizes: optional dict {level: pt}. When provided, overrides the
    master font size for that level. Only applies to body/bullet paragraphs
    (i.e. all paragraphs written here — title uses set_text_lines instead).
    Levels with no entry inherit from the master template (no override).
    """
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame

    if not append:
        tf.clear()

    if not paras:
        return

    if append and blank_before_append:
        sep = tf.add_paragraph()
        sep.add_run().text = ""
        set_bullet_none(sep)
        sep.level = 0

    for idx, ps in enumerate(paras):
        if not append and idx == 0:
            p = tf.paragraphs[0]
            for child in list(p._p):
                local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if local == "r":
                    p._p.remove(child)
        else:
            p = tf.add_paragraph()

        if ps.runs:
            for run_spec in ps.runs:
                _write_run(p, run_spec)
        else:
            p.add_run().text = ""

        if ps.bullet:
            clear_bullet_override(p)
        else:
            set_bullet_none(p)
        p.level = ps.level

        # Apply font size override for this level if specified
        if font_sizes:
            pt = font_sizes.get(ps.level)
            if pt is not None:
                try:
                    p.font.size = Pt(pt)
                except Exception as e:
                    logger.debug(f"font size apply failed (level={ps.level}): {e}")

# -------------------------
# Slide building
# -------------------------
def add_code_block_shape(slide, code_text: str) -> None:
    """Add a code-block styled textbox to the slide."""
    textbox = slide.shapes.add_textbox(
        CODE_BLOCK_LEFT, CODE_BLOCK_TOP, CODE_BLOCK_WIDTH, CODE_BLOCK_HEIGHT
    )
    fill = textbox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    line = textbox.line
    line.color.rgb = RGBColor(200, 200, 200)
    line.width = Pt(1)

    tf = textbox.text_frame
    tf.word_wrap = True
    tf.margin_bottom = Inches(0.1)
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = code_text
    p.font.name = "Courier New"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(50, 50, 50)
    if p.runs:
        set_ea_font(p.runs[0], "MS Gothic")
    send_to_back(textbox)

def insert_table_at_shape(slide, shp, lines: list[str]) -> None:
    """Insert a Markdown table into the given shape.
    Column widths are proportional to the dash counts in the separator row.
    Column alignment (left/center/right) is applied to each cell paragraph.
    """
    from pptx.enum.text import PP_ALIGN
    parsed = parse_markdown_table(lines)
    if not parsed:
        return
    h, r, col_widths, col_aligns = parsed
    n_cols = len(h)
    n_rows = len(r) + 1
    ts = slide.shapes.add_table(n_rows, n_cols, shp.left, shp.top, shp.width, shp.height)
    tbl = ts.table

    # Apply proportional column widths (last column absorbs rounding remainder)
    total_w = sum(col_widths)
    assigned = 0
    for ci, w in enumerate(col_widths[:n_cols]):
        if ci == n_cols - 1:
            tbl.columns[ci].width = shp.width - assigned
        else:
            cw = int(shp.width * w / total_w)
            tbl.columns[ci].width = cw
            assigned += cw

    # Alignment map
    ALIGN_MAP = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }

    def _set_cell(cell, text: str, align: str) -> None:
        cell.text = text
        for para in cell.text_frame.paragraphs:
            para.alignment = ALIGN_MAP.get(align, PP_ALIGN.LEFT)

    for ci, v in enumerate(h):
        _set_cell(tbl.cell(0, ci), v, col_aligns[ci] if ci < len(col_aligns) else "left")
    for ri, row in enumerate(r, 1):
        for ci, v in enumerate(row):
            if ci < n_cols:
                _set_cell(tbl.cell(ri, ci), v, col_aligns[ci] if ci < len(col_aligns) else "left")
    bring_to_front(ts)

def insert_image_at_shape(
    slide, shp, blk: list[str], md_dir: Path, ad: Path,
    placeholder_name: str = ""
) -> None:
    """Insert image with contain-fit."""
    line = blk[0].strip() if blk else ""
    m = RE_IMAGE.search(line)
    if not m or not PILImage:
        return
    alt = (m.group(1) or "").strip()
    rel = m.group(2)
    target = None
    for b in [md_dir, ad, Path.cwd()]:
        p = (b / rel).resolve()
        if p.exists():
            target = p
            break
    if not target:
        logger.debug(f"Image not found: {rel}")
        return
    try:
        with PILImage.open(target) as im:
            w, h = im.size
        iar = w / h
        bar = shp.width / shp.height
        if iar >= bar:
            nw, nh = shp.width, int(shp.width / iar)
        else:
            nw, nh = int(shp.height * iar), shp.height
        slide.shapes.add_picture(
            str(target),
            shp.left + (shp.width - nw) // 2,
            shp.top + (shp.height - nh) // 2,
            width=nw, height=nh,
        )
        logger.debug(f"  [image] inserted '{target}' alt={alt!r}")
        if alt and placeholder_name:
            cap_name = f"{placeholder_name}_caption"
            cap_shp = find_shape_by_name(slide, cap_name)
            if cap_shp is not None and getattr(cap_shp, "has_text_frame", False):
                set_text_lines(cap_shp, [alt])
                logger.debug(f"  [caption] '{cap_name}' <- {alt!r}")
            else:
                logger.debug(f"  [caption] '{cap_name}' not found, skipping")
    except Exception as e:
        logger.debug(f"Failed to insert image '{target}': {e}")

def resolve_layout_strategy(
    page_lines: list[str], front: dict, prs: Presentation, page_no: int
) -> tuple[str, dict[int, float]]:
    """Determine layout name and layout-level font_sizes for a page.
    Returns (layout_name, layout_font_sizes).
    layout_font_sizes: parsed from font_size_lN attributes in <!-- layout="..." --> tag.
    """
    forced_layout = None
    layout_font_sizes: dict[int, float] = {}

    for line in page_lines:
        m = RE_NEW_PAGE.match(line.strip())
        if m and m.group(1):
            forced_layout = m.group(1).strip()
            break

    for i, line in enumerate(page_lines):
        m = RE_LAYOUT.search(line)
        if m:
            potential_layout = m.group(1).strip()
            next_line = ""
            for j in range(i + 1, len(page_lines)):
                if page_lines[j].strip() and not RE_NEW_PAGE.match(page_lines[j].strip()):
                    next_line = page_lines[j].strip()
                    break
            h_m = RE_HEADING.match(next_line)
            if h_m and len(h_m.group(1)) <= HEADING_LEVEL_PAGE_BREAK_THRESHOLD:
                forced_layout = potential_layout
                layout_font_sizes = parse_font_sizes_from_comment(line)
                if layout_font_sizes:
                    logger.debug(f"[page {page_no}] layout font_sizes: {layout_font_sizes}")
            else:
                logger.info(f"Slide {page_no}: layout tag ignored (no heading follows).")
            break

    if forced_layout:
        if find_layout_by_name(prs, forced_layout) is not None:
            return forced_layout, layout_font_sizes
        logger.warning(
            f"Slide {page_no}: layout '{forced_layout}' not found. Falling back to auto."
        )

    heading_level = None
    for ln in page_lines:
        m = RE_HEADING.match(ln.strip())
        if m:
            heading_level = len(m.group(1))
            break

    if heading_level == 1 or bool(front):
        return "Title Slide", layout_font_sizes
    if heading_level == 2:
        return "Section Header", layout_font_sizes
    return "Title and Content", layout_font_sizes

def build_slide_from_page(prs, page_lines, md_dir, app_dir, state, page_no):
    front, body = extract_front_matter(page_lines)
    if page_no == 1:
        state["use_toc"] = front.get("toc", False)
        state["toc_title"] = front.get("toc_title", "Table of Contents")
        state["indent"] = int(front.get("indent", DEFAULT_INDENT_SPACES))
        # Collect font_size_l0 .. font_size_l4 from front matter
        font_sizes: dict[int, float] = {}
        for lv in range(5):
            key = f"font_size_l{lv}"
            if key in front:
                font_sizes[lv] = front[key]
                logger.debug(f"font_size_l{lv} = {front[key]} pt")
        state["font_sizes"] = font_sizes

    layout_name, layout_font_sizes = resolve_layout_strategy(body, front, prs, page_no)

    # Merge font_sizes: YAML(base) <- layout override
    page_font_sizes: dict[int, float] = {**state["font_sizes"], **layout_font_sizes}
    if layout_font_sizes:
        logger.debug(f"[page {page_no}] effective font_sizes after layout merge: {page_font_sizes}")

    # FIX #1: multi-level fallback — never crash due to missing layout
    layout = find_layout_by_name(prs, layout_name)
    if layout is None:
        logger.warning(
            f"Slide {page_no}: layout '{layout_name}' not found, "
            f"trying 'Title and Content'."
        )
        layout = find_layout_by_name(prs, "Title and Content")
    if layout is None:
        # Last resort: first available layout in the template
        layout = prs.slide_masters[0].slide_layouts[0]
        logger.warning(
            f"Slide {page_no}: 'Title and Content' not found either. "
            f"Using first available layout."
        )

    slide = prs.slides.add_slide(layout)
    title_shp = find_title_shape(slide)

    logger.debug(f"[page {page_no}] layout={layout_name!r}, front={bool(front)}")
    logger.debug(f"[page {page_no}] created slide: slide_layout.name={slide.slide_layout.name!r}")
    dump_slide_shapes(slide, page_no)

    body_clean = [
        ln for ln in body
        if not RE_LAYOUT.search(ln) and not RE_NEW_PAGE.match(ln.strip())
    ]
    heading_level, heading_text, heading_raw = None, None, None
    for ln in body_clean:
        m = RE_HEADING.match(ln.strip())
        if m:
            heading_level = len(m.group(1))
            heading_text = m.group(2).strip()
            heading_raw = ln.strip()
            if heading_level in (2, 3):
                state["toc"].append(TOCEntry(
                    text=heading_text, level=heading_level, target_slide=slide
                ))
            break

    logger.debug(f"[page {page_no}] heading=({heading_level}) {heading_text!r}")

    # blocks: {name: [(captured_lines, ph_font_sizes), ...]}
    blocks: dict[str, list[tuple[list[str], dict[int, float]]]] = {}
    rescue, idx = [], 0
    while idx < len(body_clean):
        m = RE_PLACEHOLDER.search(body_clean[idx])
        if m:
            name = m.group(1).strip()
            ph_font_sizes = parse_font_sizes_from_comment(body_clean[idx])
            if ph_font_sizes:
                logger.debug(f'  [page {page_no}] placeholder "{name}" font_sizes: {ph_font_sizes}')
            idx += 1
            captured = []
            if idx < len(body_clean):
                if RE_IMAGE.search(body_clean[idx].strip()):
                    captured.append(body_clean[idx])
                    idx += 1
                    while idx < len(body_clean) and RE_IMAGE.search(body_clean[idx].strip()):
                        logger.warning(
                            f"Slide {page_no}: '{name}' image limit reached. Ignoring extra."
                        )
                        idx += 1
                else:
                    while idx < len(body_clean) and body_clean[idx].strip():
                        if RE_PLACEHOLDER.search(body_clean[idx]):
                            break
                        captured.append(body_clean[idx])
                        idx += 1
            blocks.setdefault(name, []).append((captured, ph_font_sizes))
            continue
        rescue.append(body_clean[idx])
        idx += 1

    logger.debug(f"[page {page_no}] placeholder blocks: { {k: len(v) for k, v in blocks.items()} }")

    indent = state["indent"]
    verbose = state["verbose"]

    if layout_name == "Title Slide":
        info = parse_title_page_info(front, body_clean)
        if title_shp:
            set_text_lines(title_shp, [info.get("title", "")])
        for shp in slide.placeholders:
            try:
                if shp.placeholder_format and shp.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                    set_text_lines(
                        shp,
                        [ln for ln in [info.get("subtitle"), info.get("author")] if ln],
                    )
                    break
            except (AttributeError, KeyError) as e:
                logger.debug(f"Subtitle shape check failed: {e}")
    else:
        if title_shp and heading_text and heading_level in (2, 3):
            set_text_lines(title_shp, [heading_text])
            rescue = [ln for ln in rescue if ln.strip() != heading_raw]

        explicit_used = set()
        for name, blks in blocks.items():
            shp = find_shape_by_name(slide, name)
            if not shp:
                logger.debug(f'[page {page_no}] placeholder "{name}" NOT FOUND -> rescuing content')
                for blk, _ in blks:
                    rescue.extend(blk)
                continue
            explicit_used.add(id(shp))
            logger.debug(
                f'[page {page_no}] placeholder "{name}" found: '
                f'actual_name={getattr(shp, "name", None)!r}, '
                f'is_placeholder={getattr(shp, "is_placeholder", False)}, '
                f'has_text_frame={getattr(shp, "has_text_frame", False)}'
            )
            for bi, (blk, ph_font_sizes) in enumerate(blks):
                if not blk:
                    continue
                # Merge: YAML <- layout <- placeholder
                effective_font_sizes = {**page_font_sizes, **ph_font_sizes}
                if RE_IMAGE.search(blk[0].strip()):
                    logger.debug(f'  [image] block #{bi+1} -> "{name}"')
                    insert_image_at_shape(slide, shp, blk, md_dir, app_dir,
                                          placeholder_name=name)
                elif parse_markdown_table(blk):
                    logger.debug(f'  [table] block #{bi+1} -> "{name}"')
                    insert_table_at_shape(slide, shp, blk)
                elif getattr(shp, "has_text_frame", False):
                    paras, codes = build_paragraphs_from_lines(blk, indent, verbose, page_no)
                    logger.debug(f'  [text] block #{bi+1} -> "{name}" paras={len(paras)} font_sizes={effective_font_sizes or None}')
                    for c in codes:
                        add_code_block_shape(slide, c)
                    key = (id(slide), name)
                    already = state["text_written"].get(key, False)
                    write_paragraphs_to_shape(shp, paras, append=already, blank_before_append=already,
                                              font_sizes=effective_font_sizes)
                    state["text_written"][key] = True

        while rescue and not rescue[0].strip():
            rescue.pop(0)
        if any(ln.strip() for ln in rescue):
            body_shp = find_body_text_shape_excluding(slide, explicit_used)
            if body_shp:
                paras, codes = build_paragraphs_from_lines(rescue, indent, verbose, page_no)
                if blocks:
                    logger.debug(f"[page {page_no}] rescue -> body paras={len(paras)} (some content unmatched)")
                else:
                    logger.debug(f"[page {page_no}] rescue -> body paras={len(paras)} (no placeholder blocks)")
                for c in codes:
                    add_code_block_shape(slide, c)
                key = (id(slide), "__BODY__")
                # FIX #5: rescue always appends — never overwrites existing shape content
                already = state["text_written"].get(key, False)
                write_paragraphs_to_shape(body_shp, paras, append=already, blank_before_append=already,
                                          font_sizes=page_font_sizes)
                state["text_written"][key] = True
            else:
                logger.debug(f"[page {page_no}] rescue skipped (no body shape found)")

def generate_toc_slide(prs, toc_entries: list[TOCEntry], toc_title: str = "Table of Contents") -> None:
    layout = find_layout_by_name(prs, "Table of Contents")
    if layout is None:
        layout = find_layout_by_name(prs, "Title and Content")
    if layout is None:
        layout = prs.slide_masters[0].slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    title_shp = find_title_shape(slide)
    if title_shp:
        set_text_lines(title_shp, [toc_title])
    body_shp = find_body_text_shape_excluding(slide, {id(title_shp)})
    if not body_shp:
        return
    tf = body_shp.text_frame
    tf.clear()
    for i, entry in enumerate(toc_entries):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        for child in list(p._p):
            local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if local == "r":
                p._p.remove(child)
        set_bullet_none(p)
        p.level = entry.level - 2
        run = p.add_run()
        run.text = entry.text
        add_slide_hyperlink(run, slide, entry.target_slide)

def main() -> int:
    parser = argparse.ArgumentParser(description="Markdown to PPTX Injector")
    parser.add_argument("src")
    parser.add_argument("dst")
    parser.add_argument("--template", default="template.pptx")
    parser.add_argument("-v", "--verbose", action="store_true")
    args = parser.parse_args()
    setup_logging(args.verbose)
    ad = get_app_dir()

    src = resolve_path(ad, args.src)
    dst = resolve_path(ad, args.dst)
    tmp = resolve_path(ad, args.template)

    if not src.exists() or not tmp.exists():
        logger.error("File not found.")
        return 2

    prs = Presentation(str(tmp))
    state = {
        "text_written": {},
        "verbose": args.verbose,
        "toc": [],
        "use_toc": False,
        "toc_title": "Table of Contents",
        "indent": DEFAULT_INDENT_SPACES,
        "font_sizes": {},  # {level(int): pt(float)} — only specified levels
    }
    pages = split_pages(src.read_text(encoding="utf-8"))
    logger.debug(f"Processing {len(pages)} page(s)")
    for i, pg in enumerate(pages, 1):
        if any(l.strip() for l in pg):
            build_slide_from_page(prs, pg, src.parent, ad, state, i)

    if state["use_toc"] and state["toc"]:
        generate_toc_slide(prs, state["toc"], toc_title=state["toc_title"])

    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))
    logger.info(f"Done: {dst}")
    return 0

if __name__ == "__main__":
    sys.exit(main())
